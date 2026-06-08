import os
import json
import glob
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def set_shape_color(shape, fill_hex, line_hex=None, line_width_pt=1):
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = hex_to_rgb(line_hex)
        shape.line.width = Pt(line_width_pt)
    else:
        # borderless
        shape.line.fill.background()

def add_text_box(slide, left, top, width, height, text, font_name, font_size_pt, font_color_hex, bold=False, italic=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    
    tf.text = text
    
    for p in tf.paragraphs:
        p.alignment = align
        p.line_spacing = 1.2
        p.font.name = font_name
        p.font.size = Pt(font_size_pt)
        p.font.color.rgb = hex_to_rgb(font_color_hex)
        p.font.bold = bold
        p.font.italic = italic
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)
            run.font.color.rgb = hex_to_rgb(font_color_hex)
            run.font.bold = bold
            run.font.italic = italic
    return txBox

def add_header(slide, title, subtitle, font_family, color_dark, color_herbal):
    # Slide Title
    add_text_box(
        slide=slide,
        left=Inches(40 / 102.4),
        top=Inches(40 / 102.4),
        width=Inches((1024 - 80) / 102.4),
        height=Inches(45 / 102.4),
        text=title,
        font_name=font_family,
        font_size_pt=26,
        font_color_hex=color_dark,
        bold=True
    )
    # Slide Subtitle
    add_text_box(
        slide=slide,
        left=Inches(40 / 102.4),
        top=Inches(90 / 102.4),
        width=Inches((1024 - 80) / 102.4),
        height=Inches(30 / 102.4),
        text=subtitle,
        font_name=font_family,
        font_size_pt=15,
        font_color_hex=color_herbal,
        bold=False
    )

def main():
    import argparse
    parser = argparse.ArgumentParser(description="Render PPTX deck from json specs")
    parser.add_argument("--deck", type=str, default="deck_retinoids_v2", help="Name of the deck directory and files prefix")
    args = parser.parse_args()
    
    deck_name = args.deck
    
    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_root = os.path.abspath(os.path.join(script_dir, "..", ".."))
    
    tokens_path = os.path.join(project_root, "04_design_system", "design-tokens.json")
    specs_dir = os.path.join(project_root, "05_content", "specs", deck_name)
    out_dir = os.path.join(project_root, "06_render", "out")
    os.makedirs(out_dir, exist_ok=True)
    
    with open(tokens_path, "r", encoding="utf-8") as f:
        tokens = json.load(f)
        
    color = tokens["color"]
    font = tokens["font"]
    radius = tokens["radius"]
    type_scale = tokens["type"]
    
    font_family = font["family"]
    
    prs = Presentation()
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(5.625)
    
    # slide layout (use a blank layout, usually index 6 in standard templates)
    blank_layout = prs.slide_layouts[6]
    
    # Read specs
    spec_files = sorted(glob.glob(os.path.join(specs_dir, f"{deck_name}-s*.json")))
    
    for sf in spec_files:
        with open(sf, "r", encoding="utf-8") as f:
            slide_spec = json.load(f)
            
        layout = slide_spec.get("layout", "summary")
        title = slide_spec.get("title", "")
        subtitle = slide_spec.get("subtitle", "")
        body = slide_spec.get("body", [])
        media = slide_spec.get("media", {})
        components = slide_spec.get("components", {})
        disclaimers = slide_spec.get("disclaimers", [])
        source_refs = slide_spec.get("source_refs", [])
        
        slide = prs.slides.add_slide(blank_layout)
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        bg_color = color["bgAlt"] if layout in ("cover", "section_divider") else color["bg"]
        fill.fore_color.rgb = hex_to_rgb(bg_color)
        
        # Bottom disclaimers
        if disclaimers:
            disc_text = " | ".join(disclaimers)
            add_text_box(
                slide=slide,
                left=Inches(40 / 102.4),
                top=Inches(520 / 102.4),
                width=Inches(5.0),
                height=Inches(0.3),
                text=disc_text.upper(),
                font_name=font_family,
                font_size_pt=11,
                font_color_hex=color["herbal"],
                bold=True
            )
            
        # Bottom sources
        if source_refs:
            src_text = f"Источники: {', '.join(source_refs)}"
            add_text_box(
                slide=slide,
                left=Inches(5.5),
                top=Inches(520 / 102.4),
                width=Inches(4.1),
                height=Inches(0.3),
                text=src_text,
                font_name=font_family,
                font_size_pt=11,
                font_color_hex=color["text"],
                bold=False,
                align=PP_ALIGN.RIGHT
            )
            
        # Layouts Rendering
        if layout == "cover":
            # Tag
            add_text_box(
                slide=slide,
                left=Inches(40 / 102.4),
                top=Inches(180 / 102.4),
                width=Inches(4.5),
                height=Inches(0.3),
                text=slide_spec.get('section', '').upper(),
                font_name=font_family,
                font_size_pt=10,
                font_color_hex=color["herbal"],
                bold=True
            )
            # Display Title
            add_text_box(
                slide=slide,
                left=Inches(40 / 102.4),
                top=Inches(210 / 102.4),
                width=Inches(4.5),
                height=Inches(1.2),
                text=title,
                font_name=font_family,
                font_size_pt=36,
                font_color_hex=color["dark"],
                bold=True
            )
            # Subtitle
            add_text_box(
                slide=slide,
                left=Inches(40 / 102.4),
                top=Inches(340 / 102.4),
                width=Inches(4.5),
                height=Inches(1.2),
                text=subtitle,
                font_name=font_family,
                font_size_pt=15,
                font_color_hex=color["text"],
                bold=False
            )
            # Media box
            media_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(540 / 102.4), Inches(120 / 102.4),
                Inches(444 / 102.4), Inches(320 / 102.4)
            )
            set_shape_color(media_shape, color["sage"], color["border"], 1)
            # Text inside media shape
            tf = media_shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = media.get('caption', 'YM PROSKIN')
            p.font.name = font_family
            p.font.size = Pt(13)
            p.font.color.rgb = hex_to_rgb(color["herbal"])
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = font_family
                run.font.size = Pt(13)
                run.font.color.rgb = hex_to_rgb(color["herbal"])
                run.font.bold = True
            
        elif layout == "section_divider":
            # Tag
            add_text_box(
                slide=slide,
                left=Inches(40 / 102.4),
                top=Inches(180 / 102.4),
                width=Inches(9.22),
                height=Inches(0.3),
                text=slide_spec.get('section', '').upper(),
                font_name=font_family,
                font_size_pt=10,
                font_color_hex=color["herbal"],
                bold=True
            )
            # Title
            add_text_box(
                slide=slide,
                left=Inches(40 / 102.4),
                top=Inches(220 / 102.4),
                width=Inches(9.22),
                height=Inches(0.8),
                text=title,
                font_name=font_family,
                font_size_pt=36,
                font_color_hex=color["dark"],
                bold=True
            )
            # Subtitle
            add_text_box(
                slide=slide,
                left=Inches(40 / 102.4),
                top=Inches(310 / 102.4),
                width=Inches(9.22),
                height=Inches(0.8),
                text=subtitle,
                font_name=font_family,
                font_size_pt=20,
                font_color_hex=color["text"],
                bold=False
            )
            
        elif layout == "quote_callout":
            add_header(slide, title, subtitle, font_family, color["dark"], color["herbal"])
            # Quotes mark symbol
            add_text_box(
                slide=slide,
                left=Inches(4.5),
                top=Inches(150 / 102.4),
                width=Inches(1.0),
                height=Inches(0.6),
                text="“",
                font_name=font_family,
                font_size_pt=72,
                font_color_hex=color["sage"],
                bold=True,
                align=PP_ALIGN.CENTER
            )
            # Quote Text
            quote_text = ""
            for item in body:
                if item.get("type") == "quote":
                    quote_text = item.get("text", "")
            add_text_box(
                slide=slide,
                left=Inches(1.0),
                top=Inches(230 / 102.4),
                width=Inches(8.0),
                height=Inches(2.0),
                text=quote_text,
                font_name=font_family,
                font_size_pt=20,
                font_color_hex=color["dark"],
                bold=False,
                italic=True,
                align=PP_ALIGN.CENTER
            )
            
        elif layout == "two_columns_image":
            add_header(slide, title, subtitle, font_family, color["dark"], color["herbal"])
            # Bullets
            bullets_text = ""
            for item in body:
                if item.get("type") == "bullet":
                    bullets_text += f"• {item.get('text', '')}\n\n"
            bullets_text = bullets_text.strip()
            
            add_text_box(
                slide=slide,
                left=Inches(40 / 102.4),
                top=Inches(150 / 102.4),
                width=Inches(5.0),
                height=Inches(3.2),
                text=bullets_text,
                font_name=font_family,
                font_size_pt=15,
                font_color_hex=color["text"],
                bold=False
            )
            # Image Media Box
            media_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(580 / 102.4), Inches(150 / 102.4),
                Inches(404 / 102.4), Inches(320 / 102.4)
            )
            set_shape_color(media_shape, color["sage"], color["border"], 1)
            tf = media_shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = media.get('caption', '')
            p.font.name = font_family
            p.font.size = Pt(13)
            p.font.color.rgb = hex_to_rgb(color["herbal"])
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = font_family
                run.font.size = Pt(13)
                run.font.color.rgb = hex_to_rgb(color["herbal"])
                run.font.bold = True
            
        elif layout == "three_cards":
            add_header(slide, title, subtitle, font_family, color["dark"], color["herbal"])
            card_width = Inches(2.87)
            card_height = Inches(3.1)
            left_coords = [Inches(40 / 102.4), Inches(365 / 102.4), Inches(690 / 102.4)]
            
            card_idx = 0
            for item in body:
                if item.get("type") == "card" and card_idx < 3:
                    card_left = left_coords[card_idx]
                    card_idx += 1
                    
                    # Draw Card shape
                    card_shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        card_left, Inches(150 / 102.4),
                        card_width, card_height
                    )
                    set_shape_color(card_shape, color["bgAlt"], color["border"], 1)
                    
                    # Add Card Title
                    add_text_box(
                        slide=slide,
                        left=card_left + Inches(0.2),
                        top=Inches(170 / 102.4),
                        width=card_width - Inches(0.4),
                        height=Inches(0.5),
                        text=item.get("title", ""),
                        font_name=font_family,
                        font_size_pt=18,
                        font_color_hex=color["dark"],
                        bold=True
                    )
                    # Add Card Text
                    add_text_box(
                        slide=slide,
                        left=card_left + Inches(0.2),
                        top=Inches(230 / 102.4),
                        width=card_width - Inches(0.4),
                        height=Inches(2.1),
                        text=item.get("text", ""),
                        font_name=font_family,
                        font_size_pt=13,
                        font_color_hex=color["text"]
                    )
                    
        elif layout == "timeline_protocol":
            add_header(slide, title, subtitle, font_family, color["dark"], color["herbal"])
            card_width = Inches(2.05)
            card_height = Inches(3.1)
            left_coords = [
                Inches(40 / 102.4),
                Inches(280 / 102.4),
                Inches(520 / 102.4),
                Inches(760 / 102.4)
            ]
            
            step_idx = 0
            for item in body:
                if item.get("type") == "step" and step_idx < 4:
                    item_left = left_coords[step_idx]
                    step_idx += 1
                    
                    # Draw box
                    box_shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        item_left, Inches(150 / 102.4),
                        card_width, card_height
                    )
                    set_shape_color(box_shape, color["bgAlt"], color["border"], 1)
                    
                    # Draw step number badge (as shape)
                    badge_shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        item_left + Inches(0.15), Inches(170 / 102.4),
                        Inches(0.6), Inches(0.3)
                    )
                    set_shape_color(badge_shape, color["herbal"])
                    badge_tf = badge_shape.text_frame
                    p = badge_tf.paragraphs[0]
                    p.text = str(step_idx)
                    p.font.name = font_family
                    p.font.size = Pt(10)
                    p.font.bold = True
                    p.font.color.rgb = hex_to_rgb(color["bg"])
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.name = font_family
                        run.font.size = Pt(10)
                        run.font.bold = True
                        run.font.color.rgb = hex_to_rgb(color["bg"])
                    
                    # Add Step Label
                    add_text_box(
                        slide=slide,
                        left=item_left + Inches(0.15),
                        top=Inches(215 / 102.4),
                        width=card_width - Inches(0.3),
                        height=Inches(0.4),
                        text=item.get("label", ""),
                        font_name=font_family,
                        font_size_pt=14,
                        font_color_hex=color["dark"],
                        bold=True
                    )
                    # Add Step Desc
                    add_text_box(
                        slide=slide,
                        left=item_left + Inches(0.15),
                        top=Inches(265 / 102.4),
                        width=card_width - Inches(0.3),
                        height=Inches(1.8),
                        text=item.get("desc", ""),
                        font_name=font_family,
                        font_size_pt=11,
                        font_color_hex=color["text"]
                    )
                    
        elif layout == "comparison":
            add_header(slide, title, subtitle, font_family, color["dark"], color["herbal"])
            card_width = Inches(4.41)
            card_height = Inches(3.1)
            
            left_label, left_text = "", ""
            right_label, right_text = "", ""
            for item in body:
                if item.get("type") == "compare_left":
                    left_label = item.get("label", "")
                    left_text = item.get("text", "")
                elif item.get("type") == "compare_right":
                    right_label = item.get("label", "")
                    right_text = item.get("text", "")
            
            # Left Card shape
            left_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(40 / 102.4), Inches(150 / 102.4),
                card_width, card_height
            )
            set_shape_color(left_shape, color["bgAlt"], color["border"], 1)
            add_text_box(
                slide=slide,
                left=Inches(60 / 102.4),
                top=Inches(170 / 102.4),
                width=card_width - Inches(0.4),
                height=Inches(0.4),
                text=left_label,
                font_name=font_family,
                font_size_pt=18,
                font_color_hex=color["herbal"],
                bold=True
            )
            add_text_box(
                slide=slide,
                left=Inches(60 / 102.4),
                top=Inches(220 / 102.4),
                width=card_width - Inches(0.4),
                height=Inches(2.2),
                text=left_text,
                font_name=font_family,
                font_size_pt=15,
                font_color_hex=color["text"]
            )
            
            # Right Card shape
            right_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(532 / 102.4), Inches(150 / 102.4),
                card_width, card_height
            )
            set_shape_color(right_shape, color["bgAlt"], color["border"], 1)
            add_text_box(
                slide=slide,
                left=Inches(552 / 102.4),
                top=Inches(170 / 102.4),
                width=card_width - Inches(0.4),
                height=Inches(0.4),
                text=right_label,
                font_name=font_family,
                font_size_pt=18,
                font_color_hex=color["herbal"],
                bold=True
            )
            add_text_box(
                slide=slide,
                left=Inches(552 / 102.4),
                top=Inches(220 / 102.4),
                width=card_width - Inches(0.4),
                height=Inches(2.2),
                text=right_text,
                font_name=font_family,
                font_size_pt=15,
                font_color_hex=color["text"]
            )
            
        elif layout == "statistics":
            add_header(slide, title, subtitle, font_family, color["dark"], color["herbal"])
            card_width = Inches(2.87)
            card_height = Inches(3.1)
            left_coords = [Inches(40 / 102.4), Inches(365 / 102.4), Inches(690 / 102.4)]
            
            stat_idx = 0
            for item in body:
                if item.get("type") == "stat_item" and stat_idx < 3:
                    card_left = left_coords[stat_idx]
                    stat_idx += 1
                    
                    # Draw Stat shape
                    stat_shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        card_left, Inches(150 / 102.4),
                        card_width, card_height
                    )
                    set_shape_color(stat_shape, color["bgAlt"], color["border"], 1)
                    
                    # Stat Number
                    add_text_box(
                        slide=slide,
                        left=card_left + Inches(0.2),
                        top=Inches(165 / 102.4),
                        width=card_width - Inches(0.4),
                        height=Inches(0.8),
                        text=item.get("number", ""),
                        font_name=font_family,
                        font_size_pt=52,
                        font_color_hex=color["herbal"],
                        bold=True
                    )
                    # Stat Label
                    add_text_box(
                        slide=slide,
                        left=card_left + Inches(0.2),
                        top=Inches(250 / 102.4),
                        width=card_width - Inches(0.4),
                        height=Inches(0.4),
                        text=item.get("label", ""),
                        font_name=font_family,
                        font_size_pt=18,
                        font_color_hex=color["dark"],
                        bold=True
                    )
                    # Stat Desc
                    add_text_box(
                        slide=slide,
                        left=card_left + Inches(0.2),
                        top=Inches(295 / 102.4),
                        width=card_width - Inches(0.4),
                        height=Inches(1.4),
                        text=item.get("desc", ""),
                        font_name=font_family,
                        font_size_pt=12,
                        font_color_hex=color["text"]
                    )
                    
        elif layout == "contraindication_alert":
            add_header(slide, title, subtitle, font_family, color["dark"], color["herbal"])
            alert = components.get("alert", {})
            alert_title = alert.get("title", "Противопоказание")
            alert_text = alert.get("text", "")
            
            # Alert Box shape
            alert_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(40 / 102.4), Inches(150 / 102.4),
                Inches(9.22), Inches(2.73)
            )
            set_shape_color(alert_shape, color["bgAlt"], color["warn"], 2)
            
            # Icon
            add_text_box(
                slide=slide,
                left=Inches(0.7),
                top=Inches(225 / 102.4),
                width=Inches(0.8),
                height=Inches(0.8),
                text="⚠️",
                font_name=font_family,
                font_size_pt=48,
                font_color_hex=color["warn"],
                align=PP_ALIGN.CENTER
            )
            # Warning content
            add_text_box(
                slide=slide,
                left=Inches(1.66),
                top=Inches(180 / 102.4),
                width=Inches(7.5),
                height=Inches(0.4),
                text=alert_title,
                font_name=font_family,
                font_size_pt=20,
                font_color_hex=color["warn"],
                bold=True
            )
            add_text_box(
                slide=slide,
                left=Inches(1.66),
                top=Inches(230 / 102.4),
                width=Inches(7.5),
                height=Inches(1.8),
                text=alert_text,
                font_name=font_family,
                font_size_pt=15,
                font_color_hex=color["text"],
                bold=True
            )
            
        elif layout == "summary":
            add_header(slide, title, subtitle, font_family, color["dark"], color["herbal"])
            bullets_text = ""
            for item in body:
                if item.get("type") == "bullet":
                    bullets_text += f"• {item.get('text', '')}\n\n"
            bullets_text = bullets_text.strip()
            
            # Card Box
            summary_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(40 / 102.4), Inches(146 / 102.4),
                Inches(9.22), Inches(3.125)
            )
            set_shape_color(summary_shape, color["bgAlt"], color["border"], 1)
            
            add_text_box(
                slide=slide,
                left=Inches(0.7),
                top=Inches(177 / 102.4),
                width=Inches(8.6),
                height=Inches(2.5),
                text=bullets_text,
                font_name=font_family,
                font_size_pt=15,
                font_color_hex=color["text"],
                bold=False
            )
            
    out_pptx_path = os.path.join(out_dir, f"{deck_name}.pptx")
    prs.save(out_pptx_path)
    print(f"PPTX presentation written successfully to {out_pptx_path}")

if __name__ == "__main__":
    main()
