import os
import json
import glob
import subprocess
import urllib.request
import sys
from pptx import Presentation
from pptx.dml.color import RGBColor

def main():
    import argparse
    parser = argparse.ArgumentParser(description="QA audit on rendered HTML and PPTX files")
    parser.add_argument("--deck", type=str, default="deck_retinoids_v2", help="Name of the deck to audit")
    args = parser.parse_args()
    
    deck_name = args.deck
    
    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_root = os.path.abspath(os.path.join(script_dir, "..", ".."))
    
    tokens_path = os.path.join(project_root, "04_design_system", "design-tokens.json")
    specs_dir = os.path.join(project_root, "05_content", "specs", deck_name)
    logs_dir = os.path.join(project_root, "ops", "logs")
    os.makedirs(logs_dir, exist_ok=True)
    
    # 1. Download Arimo font if not present
    fonts_dir = os.path.join(project_root, "04_design_system", "fonts")
    os.makedirs(fonts_dir, exist_ok=True)
    font_path = os.path.join(fonts_dir, "Arimo[wght].ttf")
    if not os.path.exists(font_path):
        print("Local Arimo font file not found. Downloading...")
        url = "https://github.com/google/fonts/raw/main/ofl/arimo/Arimo%5Bwght%5D.ttf"
        try:
            req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
            with urllib.request.urlopen(req, timeout=30) as r:
                data = r.read()
            with open(font_path, "wb") as f:
                f.write(data)
            print(f"Font downloaded successfully to {font_path} ({len(data)} bytes)")
        except Exception as e:
            print(f"Warning: Failed to download local font: {e}")
            
    # 2. Load and verify design tokens
    with open(tokens_path, "r", encoding="utf-8") as f:
        tokens = json.load(f)
        
    pdf_errors = []
    pptx_errors = []
    warnings = []
    
    # Load asset provenance registry
    registry_path = os.path.join(project_root, "04_design_system", "assets", "asset_provenance.json")
    registry = {}
    if os.path.exists(registry_path):
        try:
            with open(registry_path, "r", encoding="utf-8") as rf:
                registry = json.load(rf)
        except Exception as e:
            pdf_errors.append(f"Provenance registry load error: {e}")
            pptx_errors.append(f"Provenance registry load error: {e}")

    # Load HTML output content for attribution checking
    html_path = os.path.join(project_root, "06_render", "out", f"{deck_name}.html")
    html_content = ""
    if os.path.exists(html_path):
        try:
            with open(html_path, "r", encoding="utf-8") as hf:
                html_content = hf.read()
        except Exception as e:
            pdf_errors.append(f"Error reading HTML file: {e}")
    else:
        pdf_errors.append(f"HTML presentation file not found: {html_path}")
    
    # Check Zero Black Policy
    color_tokens = tokens.get("color", {})
    for col_name, col_val in color_tokens.items():
        val_clean = str(col_val).strip().lower()
        if val_clean in ("#000000", "#000", "black", "rgb(0,0,0)", "rgba(0,0,0,1)"):
            pdf_errors.append(f"Zero Black violation: token 'color.{col_name}' is set to black ({col_val})")
            
    # Check Font Family
    font_family = tokens.get("font", {}).get("family", "")
    if font_family != "Arimo":
        pdf_errors.append(f"Font Family check failed: expected 'Arimo', got '{font_family}'")
        
    # Load safety config
    safety_config_path = os.path.join(project_root, "05_content", "safety_config.json")
    pregnancy_required_decks = []
    if os.path.exists(safety_config_path):
        try:
            with open(safety_config_path, "r", encoding="utf-8") as scf:
                safety_config = json.load(scf)
                pregnancy_required_decks = safety_config.get("pregnancy_required_decks", [])
        except Exception as e:
            print(f"Warning: Failed to load safety config: {e}")

    topic_name = deck_name.replace("deck_", "") if deck_name.startswith("deck_") else deck_name
    
    # Load all active facts for this topic to check for pregnancy keywords
    has_pregnancy_fact_in_graph = False
    facts_dir = os.path.join(project_root, "03_knowledge_graph", "facts")
    if os.path.exists(facts_dir):
        for f_file in os.listdir(facts_dir):
            if f_file.endswith(".json"):
                try:
                    with open(os.path.join(facts_dir, f_file), "r", encoding="utf-8") as ff:
                        f_data = json.load(ff)
                    f_entity = str(f_data.get("entity_id", "")).lower()
                    if f_entity == topic_name or f_entity in topic_name or topic_name in f_entity:
                        statement = str(f_data.get("statement", "")).lower()
                        if any(kw in statement for kw in ["беременнос", "pregnancy", "лактаци", "lactation"]):
                            has_pregnancy_fact_in_graph = True
                            break
                except Exception:
                    pass

    require_pregnancy_slide = (deck_name in pregnancy_required_decks) or has_pregnancy_fact_in_graph

    # 3. Check Slide Specs (Source Refs and Pregnancy/Safety checks)
    spec_files = sorted(glob.glob(os.path.join(specs_dir, f"{deck_name}-s*.json")))
    total_slides = len(spec_files)
    
    quarantined_facts = {
        "fact_0003", "fact_0008", "fact_0013", "fact_0016",
        "fact_0023", "fact_0026", "fact_0027",
        "fact_0031", "fact_0032", "fact_0033", "fact_0037",
        "fact_0040", "fact_0042", "fact_0043"
    }
    
    has_pregnancy_slide = False
    has_safety_slide = False
    
    for sf in spec_files:
        with open(sf, "r", encoding="utf-8") as f:
            slide = json.load(f)
        
        layout = slide.get("layout", "")
        refs = slide.get("source_refs", [])
        
        # [GATE] Asset Provenance and License Verification
        media = slide.get("media", {})
        asset_name = media.get("asset", "")
        if asset_name:
            asset_data = registry.get(asset_name)
            if not asset_data:
                base_name = os.path.splitext(asset_name)[0]
                asset_data = registry.get(base_name)
                
            if not asset_data:
                err = f"GATE BLOCK: Asset '{asset_name}' used in slide '{os.path.basename(sf)}' is NOT registered in asset_provenance.json."
                pdf_errors.append(err)
                pptx_errors.append(err)
            else:
                # [GATE] SVG Bounds and Layout Check
                asset_path_rel = asset_data.get("path", "")
                if asset_path_rel and asset_path_rel.endswith(".svg"):
                    asset_path_abs = os.path.join(project_root, asset_path_rel)
                    if os.path.exists(asset_path_abs):
                        try:
                            svg_check_script = os.path.join(project_root, "ops", "scripts", "qa_svg_bounds.js")
                            res_svg = subprocess.run(
                                ["node", svg_check_script, asset_path_abs, "5"],
                                capture_output=True,
                                text=True,
                                timeout=20
                            )
                            if res_svg.returncode != 0:
                                err = f"GATE BLOCK: SVG bounds check failed for '{asset_name}' ({asset_path_rel}): {res_svg.stderr or res_svg.stdout}"
                                pdf_errors.append(err)
                                pptx_errors.append(err)
                        except Exception as e:
                            err = f"GATE BLOCK: SVG bounds check exception for '{asset_name}': {e}"
                            pdf_errors.append(err)
                            pptx_errors.append(err)

                license_type = str(asset_data.get("license", "")).lower()
                allowed_licenses = ["own/generated", "generated/own", "cc0", "cc-by", "cc-by-4.0", "public-domain"]
                
                # Check for forbidden strings (like -sa, -nc, or biorender)
                is_forbidden = any(bad in license_type for bad in ["-sa", "-nc", "biorender"]) or license_type not in allowed_licenses
                if is_forbidden:
                    err = f"GATE BLOCK: Asset '{asset_name}' has forbidden or invalid license type '{license_type}'."
                    pdf_errors.append(err)
                    pptx_errors.append(err)
                    
                # CC-BY check
                attribution = asset_data.get("attribution", "")
                if "cc-by" in license_type and attribution:
                    if html_content and attribution not in html_content:
                        err = f"GATE BLOCK: CC-BY asset '{asset_name}' attribution '{attribution}' is missing from the output HTML/PDF."
                        pdf_errors.append(err)
                        
                # Molecule check
                asset_type = asset_data.get("type", "")
                source_truth = str(asset_data.get("source_of_truth", "")).lower()
                if asset_type == "molecule":
                    import re
                    if not re.search(r"pubchem cid \d+", source_truth):
                        err = f"GATE BLOCK: Molecular asset '{asset_name}' must reference a valid PubChem CID in source_of_truth (got '{source_truth}')."
                        pdf_errors.append(err)
                        pptx_errors.append(err)
                        
                # Diagram-fact check for pilot mechanisms
                if asset_name in ["rar_rxr_mechanism.svg", "skin_layers_turnover.svg", "collagen_synthesis.svg", "ascorbic_acid_absorption.svg", "ceramide_synthesis.svg", "melanosome_transfer.svg", "desmosome_desmolysis.svg"]:
                    if asset_name == "rar_rxr_mechanism.svg":
                        if not any(f in refs for f in ["fact_0001", "fact_0002"]):
                            err = f"GATE BLOCK: Slide '{os.path.basename(sf)}' uses rar_rxr_mechanism.svg but does not cite fact_0001 or fact_0002 in source_refs."
                            pdf_errors.append(err)
                            pptx_errors.append(err)
                    elif asset_name == "skin_layers_turnover.svg":
                        if not any(f in refs for f in ["fact_0006", "fact_0047", "fact_0048"]):
                            err = f"GATE BLOCK: Slide '{os.path.basename(sf)}' uses skin_layers_turnover.svg but does not cite fact_0006, fact_0047, or fact_0048 in source_refs."
                            pdf_errors.append(err)
                            pptx_errors.append(err)
                    elif asset_name == "collagen_synthesis.svg":
                        if not any(f in refs for f in ["fact_0021", "fact_0030", "fact_0046"]):
                            err = f"GATE BLOCK: Slide '{os.path.basename(sf)}' uses collagen_synthesis.svg but does not cite fact_0021, fact_0030, or fact_0046 in source_refs."
                            pdf_errors.append(err)
                            pptx_errors.append(err)
                    elif asset_name == "ascorbic_acid_absorption.svg":
                        if not any(f in refs for f in ["fact_0028", "fact_0029", "fact_0044", "fact_0045"]):
                            err = f"GATE BLOCK: Slide '{os.path.basename(sf)}' uses ascorbic_acid_absorption.svg but does not cite fact_0028, fact_0029, fact_0044, or fact_0045 in source_refs."
                            pdf_errors.append(err)
                            pptx_errors.append(err)
                    elif asset_name == "ceramide_synthesis.svg":
                        if "fact_0024" not in refs:
                            err = f"GATE BLOCK: Slide '{os.path.basename(sf)}' uses ceramide_synthesis.svg but does not cite fact_0024 in source_refs."
                            pdf_errors.append(err)
                            pptx_errors.append(err)
                    elif asset_name == "melanosome_transfer.svg":
                        if "fact_0034" not in refs:
                            err = f"GATE BLOCK: Slide '{os.path.basename(sf)}' uses melanosome_transfer.svg but does not cite fact_0034 in source_refs."
                            pdf_errors.append(err)
                            pptx_errors.append(err)
                    elif asset_name == "desmosome_desmolysis.svg":
                        if "fact_0049" not in refs:
                            err = f"GATE BLOCK: Slide '{os.path.basename(sf)}' uses desmosome_desmolysis.svg but does not cite fact_0049 in source_refs."
                            pdf_errors.append(err)
                            pptx_errors.append(err)
                            
                # Placeholder caption check
                if "placeholder" in asset_name.lower():
                    caption = media.get("caption", "")
                    if not caption.startswith("[Placeholder]"):
                        err = f"GATE BLOCK: Placeholder asset '{asset_name}' caption '{caption}' in slide '{os.path.basename(sf)}' must start with '[Placeholder]'."
                        pdf_errors.append(err)
                        pptx_errors.append(err)
        
        if layout == "contraindication_alert":
            has_safety_slide = True
            
        # Check source refs on clinical slides (non-cover and non-section_divider)
        if layout not in ("cover", "section_divider"):
            if not refs:
                warnings.append(f"Slide '{os.path.basename(sf)}' is clinical but has no source references.")
            else:
                # [GATE] Check if any cited fact is quarantined
                for r in refs:
                    if r in quarantined_facts:
                        pdf_errors.append(f"GATE BLOCK: Slide '{os.path.basename(sf)}' cites quarantined fact '{r}'.")
                        pptx_errors.append(f"GATE BLOCK: Slide '{os.path.basename(sf)}' cites quarantined fact '{r}'.")
                        
            # Claim validation check for clinical slides
            slide_text_parts = []
            title = str(slide.get("title", ""))
            subtitle = str(slide.get("subtitle", ""))
            if title: slide_text_parts.append(title)
            if subtitle: slide_text_parts.append(subtitle)
            for item in slide.get("body", []):
                if "text" in item:
                    slide_text_parts.append(item["text"])
                if "title" in item:
                    slide_text_parts.append(item["title"])
                if "desc" in item:
                    slide_text_parts.append(item["desc"])
                if "label" in item:
                    slide_text_parts.append(item["label"])
            alert = slide.get("components", {}).get("alert", {})
            if alert:
                if "title" in alert:
                    slide_text_parts.append(alert["title"])
                if "text" in alert:
                    slide_text_parts.append(alert["text"])
            
            slide_combined_text = " ".join(slide_text_parts).lower()
            
            # Load cited facts
            cited_facts = []
            for r in refs:
                if r.startswith("fact_"):
                    fact_path = os.path.join(project_root, "03_knowledge_graph", "facts", f"{r}.json")
                    if os.path.exists(fact_path):
                        try:
                            with open(fact_path, "r", encoding="utf-8") as ff:
                                f_data = json.load(ff)
                            cited_facts.append(f_data)
                        except Exception:
                            pass
                            
            # Validate claims matching categories
            categories = [
                {
                    "name": "Pregnancy/Lactation",
                    "keywords": ["беременнос", "pregnancy", "лактаци", "lactation"]
                },
                {
                    "name": "Contraindication",
                    "keywords": ["противопоказа", "contraindicat"]
                },
                {
                    "name": "Safety/Tolerability",
                    "keywords": ["безопас", "safe"]
                }
            ]
            
            for cat in categories:
                cat_name = cat["name"]
                cat_kws = cat["keywords"]
                
                if any(kw in slide_combined_text for kw in cat_kws):
                    fact_has_support = False
                    for fact in cited_facts:
                        statement = str(fact.get("statement", "")).lower()
                        if any(kw in statement for kw in cat_kws):
                            fact_has_support = True
                            break
                            
                    # For safety, allow matches on "tolerab" / "well-tolerated" / "benefit" in English
                    if cat_name == "Safety/Tolerability" and not fact_has_support:
                        for fact in cited_facts:
                            statement = str(fact.get("statement", "")).lower()
                            if any(kw in statement for kw in ["tolerab", "well-tolerated", "benefit"]):
                                fact_has_support = True
                                break
                                
                    if not fact_has_support:
                        err_msg = f"Claim Validation failed in slide '{os.path.basename(sf)}': Slide contains '{cat_name}' clinical claim, but none of the cited facts {refs} contain supporting keywords {cat_kws}."
                        pdf_errors.append(err_msg)
                        pptx_errors.append(err_msg)

        # Check for pregnancy reference (pregnancy contraindication check)
        title = str(slide.get("title", "")).lower()
        subtitle = str(slide.get("subtitle", "")).lower()
        body_text = ""
        for item in slide.get("body", []):
            body_text += str(item.get("text", "")).lower()
        notes = str(slide.get("notes", "")).lower()
        alert = slide.get("components", {}).get("alert", {})
        alert_text = (str(alert.get("title", "")) + " " + str(alert.get("text", ""))).lower() if alert else ""
        
        combined_text = title + " " + subtitle + " " + body_text + " " + notes + " " + alert_text
        if "беременнос" in combined_text or "pregnancy" in combined_text or "лактаци" in combined_text:
            has_pregnancy_slide = True
            
    if not has_safety_slide:
        pdf_errors.append("Safety slide check failed: No safety slide (layout 'contraindication_alert') found.")
        pptx_errors.append("Safety slide check failed: No safety slide (layout 'contraindication_alert') found.")
        
    if require_pregnancy_slide and not has_pregnancy_slide:
        pdf_errors.append("Pregnancy precaution check failed: Pregnancy slide required for this topic, but none found.")
        pptx_errors.append("Pregnancy precaution check failed: Pregnancy slide required for this topic, but none found.")
        
    # 4. Run Playwright Font Load Check for PDF/HTML output
    playwright_script = os.path.join(script_dir, "qa_font_check.js")
    font_check_passed = False
    font_check_offline_passed = False
    
    try:
        res = subprocess.run(
            ["node", playwright_script, f"{deck_name}.html"],
            cwd=project_root,
            capture_output=True,
            text=True,
            timeout=30
        )
        print("Playwright stdout:\n", res.stdout)
        if res.stderr:
            print("Playwright stderr:\n", res.stderr)
            
        if "RESULT_ARIMO_LOADED=true" in res.stdout:
            font_check_passed = True
        if "RESULT_ARIMO_OFFLINE_LOADED=true" in res.stdout:
            font_check_offline_passed = True
            
        if res.returncode != 0:
            pdf_errors.append("Playwright font rendering audit returned non-zero exit code.")
    except Exception as e:
        pdf_errors.append(f"Failed to run Playwright font check: {e}")
        
    if not font_check_passed:
        pdf_errors.append("Playwright check: Arimo webfont is not loaded.")
    if not font_check_offline_passed:
        warnings.append("Offline check: Arimo font could not be loaded in offline mode. Ensure local @font-face is active.")
        
    # 4.5 Run Playwright Layout Bounds Check for PDF/HTML output
    layout_bounds_script = os.path.join(script_dir, "qa_layout_bounds.js")
    try:
        res_bounds = subprocess.run(
            ["node", layout_bounds_script, f"{deck_name}.html"],
            cwd=project_root,
            capture_output=True,
            text=True,
            timeout=30
        )
        print("Playwright layout bounds stdout:\n", res_bounds.stdout)
        if res_bounds.stderr:
            print("Playwright layout bounds stderr:\n", res_bounds.stderr)
            
        if res_bounds.returncode != 0:
            pdf_errors.append("Playwright layout/brand bounds audit failed or returned non-zero exit code.")
            try:
                errs = json.loads(res_bounds.stderr or res_bounds.stdout)
                for e in errs:
                    pdf_errors.append(f"HTML Slide {e.get('slide', 'N/A')}: {e.get('check')} - {e.get('details')}")
            except Exception:
                pass
    except Exception as e:
        pdf_errors.append(f"Failed to run Playwright layout bounds check: {e}")
        
    # 5. PPTX Specific Audits (Zero Black, Arimo, Structure)
    pptx_path = os.path.join(project_root, "06_render", "out", f"{deck_name}.pptx")
    if not os.path.exists(pptx_path):
        pptx_errors.append(f"PPTX file not found in render output: {pptx_path}")
    else:
        try:
            prs = Presentation(pptx_path)
            if len(prs.slides) != total_slides:
                pptx_errors.append(f"PPTX slide count mismatch: specs={total_slides}, pptx={len(prs.slides)}")
                
            for idx, slide in enumerate(prs.slides):
                # Load corresponding spec to resolve CC-BY attribution requirements
                with open(spec_files[idx], "r", encoding="utf-8") as sf_f:
                    spec_data = json.load(sf_f)
                
                media = spec_data.get("media", {})
                asset_name = media.get("asset", "")
                attribution = ""
                if asset_name:
                    asset_data = registry.get(asset_name) or registry.get(os.path.splitext(asset_name)[0]) or {}
                    license_type = str(asset_data.get("license", "")).lower()
                    if "cc-by" in license_type:
                        attribution = asset_data.get("attribution", "")

                slide_text = ""
                # Brand Logo presence checks in PPTX
                slide_layout = spec_data.get("layout", "summary")
                is_cover = (idx == 0)
                is_closing = (idx == total_slides - 1)
                
                has_cover_logo = False
                has_section_logo = False
                has_footer_logo = False
                has_closing_badge = False
                
                for shape in slide.shapes:
                    # shape.shape_type == 13 is PICTURE
                    if shape.shape_type == 13:
                        left_in = shape.left.inches
                        top_in = shape.top.inches
                        
                        if is_cover:
                            if abs(left_in - 40/102.4) < 0.05 and abs(top_in - 410/102.4) < 0.05:
                                has_cover_logo = True
                        elif slide_layout == "section_divider":
                            if abs(left_in - 40/102.4) < 0.05 and abs(top_in - 110/102.4) < 0.05:
                                has_section_logo = True
                        else:
                            # content footer logo (mark)
                            if abs(left_in - 40/102.4) < 0.05 and abs(top_in - 543/102.4) < 0.05:
                                has_footer_logo = True
                                
                        if is_closing and slide_layout == "summary":
                            # badge centered
                            if abs(left_in - 420/102.4) < 0.1 and abs(top_in - 196/102.4) < 0.1:
                                has_closing_badge = True

                    # Check shapes solid fills
                    if hasattr(shape, 'fill') and shape.fill.type == 1:  # solid fill
                        rgb = shape.fill.fore_color.rgb
                        if rgb == RGBColor(0, 0, 0) or rgb == RGBColor(255, 255, 255):
                            pptx_errors.append(f"Slide {idx + 1}: PPTX Zero Black violation (shape solid color is {rgb})")
                            
                    # Check text box properties
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            slide_text += paragraph.text + " "
                            for run in paragraph.runs:
                                if run.font.name != "Arimo":
                                    pptx_errors.append(f"Slide {idx + 1}: PPTX font fallback detected. Font family is '{run.font.name}', expected 'Arimo'.")
                                if run.font.color.type == 1:  # RGBColor
                                    rgb = run.font.color.rgb
                                    if rgb == RGBColor(0, 0, 0) or rgb == RGBColor(255, 255, 255):
                                        pptx_errors.append(f"Slide {idx + 1}: PPTX Zero Black violation (text run color is {rgb})")
                                        
                if is_cover and not has_cover_logo:
                    pptx_errors.append(f"Slide {idx + 1}: Cover slide is missing descriptor logo shape in PPTX")
                if slide_layout == "section_divider" and not has_section_logo:
                    pptx_errors.append(f"Slide {idx + 1}: Section divider slide is missing horizontal logo shape in PPTX")
                if slide_layout not in ("cover", "section_divider") and not has_footer_logo:
                    pptx_errors.append(f"Slide {idx + 1}: Content slide is missing footer mark logo shape in PPTX")
                if is_closing and slide_layout == "summary" and not has_closing_badge:
                    pptx_errors.append(f"Slide {idx + 1}: Closing summary slide is missing badge logo shape in PPTX")
                    
                if attribution and attribution.strip().lower() not in slide_text.strip().lower():
                    pptx_errors.append(f"Slide {idx + 1}: GATE BLOCK: CC-BY attribution '{attribution}' is missing from PPTX slide.")
        except Exception as e:
            pptx_errors.append(f"Failed to audit PPTX: {e}")
            
    # 5.5 SVG Bounds Check (internal overlaps + edge clipping via qa_svg_bounds.js)
    svg_bounds_errors = []
    mechanisms_dir = os.path.join(project_root, "04_design_system", "assets", "mechanisms")
    qa_svg_script = os.path.join(project_root, "ops", "scripts", "qa_svg_bounds.js")
    if os.path.isdir(mechanisms_dir) and os.path.isfile(qa_svg_script):
        svg_files = sorted(f for f in os.listdir(mechanisms_dir) if f.endswith(".svg"))
        for svg_fname in svg_files:
            svg_path = os.path.join(mechanisms_dir, svg_fname)
            try:
                result = subprocess.run(
                    ["node", qa_svg_script, svg_path, "5"],
                    capture_output=True, text=True, timeout=30
                )
                if result.returncode != 0:
                    import re
                    checks = re.findall(r'"check":\s*"([^"]+)"', result.stderr + result.stdout)
                    for chk in checks:
                        svg_bounds_errors.append(f"SVG {svg_fname}: {chk}")
            except Exception as e:
                svg_bounds_errors.append(f"SVG bounds check error for {svg_fname}: {e}")
    
    pdf_errors.extend(svg_bounds_errors)

    # 6. Write QA Audit Report ops/logs/qa_deck_v2.md
    pdf_status = "FAIL" if pdf_errors else "PASS"
    pptx_status = "FAIL" if pptx_errors else "PASS"
    overall_status = "FAIL" if (pdf_errors or pptx_errors) else "PASS"
    
    report_lines = [
        "# QA Audit Report — Phase 3.5: Design System & Dual Render",
        "",
        f"**Overall Audit Status**: {overall_status}",
        "",
        "## PDF / HTML Render Audit",
        f"- **Status**: {pdf_status}",
        "",
        "### Errors",
    ]
    if pdf_errors:
        for err in pdf_errors:
            report_lines.append(f"- {err}")
    else:
        report_lines.append("- No PDF rendering or citation errors found.")
        
    report_lines.extend([
        "",
        "## PPTX (PowerPoint) Render Audit",
        f"- **Status**: {pptx_status}",
        "",
        "### Errors",
    ])
    if pptx_errors:
        for err in pptx_errors:
            report_lines.append(f"- {err}")
    else:
        report_lines.append("- No PPTX font, color, or structural errors found.")
        
    report_lines.extend([
        "",
        "## Warnings",
    ])
    if warnings:
        for warn in warnings:
            report_lines.append(f"- {warn}")
    else:
        report_lines.append("- No warnings.")
        
    report_lines.extend([
        "",
        "## Checked Invariants",
        f"- Zero Black Policy: {'PASS' if not any('Zero Black' in e for e in pdf_errors + pptx_errors) else 'FAIL'}",
        f"- Arimo font defined in tokens: {'PASS' if font_family == 'Arimo' else 'FAIL'}",
        f"- Safety slide present: {'PASS' if has_safety_slide else 'FAIL'}",
        f"- Pregnancy slide present (when required): {'PASS' if (not require_pregnancy_slide or has_pregnancy_slide) else 'FAIL'}",
        f"- Playwright Font Load (Online): {'PASS' if font_check_passed else 'FAIL'}",
        f"- Playwright Font Load (Offline): {'PASS' if font_check_offline_passed else 'FAIL'}",
        f"- PPTX Element Font Check (Arimo): {'PASS' if not any('font' in e for e in pptx_errors) else 'FAIL'}",
        f"- Total slides checked: {total_slides}",
    ])
    
    report_content = "\n".join(report_lines) + "\n"
    topic_name = deck_name.replace("deck_", "") if deck_name.startswith("deck_") else deck_name
    report_path = os.path.join(logs_dir, f"qa_deck_{topic_name}.md")
    with open(report_path, "w", encoding="utf-8") as f:
        f.write(report_content)
        
    print(f"\nQA audit complete. Report written to {report_path}")
    print(f"Overall Status: {overall_status}")
    
    if overall_status == "FAIL":
        sys.exit(1)
    else:
        sys.exit(0)

if __name__ == "__main__":
    main()
